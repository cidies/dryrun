# Noch zu tun:

from flask import Flask, render_template, request, jsonify, redirect, url_for, flash, send_from_directory
from threading import Timer
import json
import os
import io
from flask_toastr import Toastr
from flask import flash, get_flashed_messages
import time
from datetime import datetime
from flask_socketio import SocketIO, emit, Namespace
from threading import Thread
import logging
import smtplib
from email.mime.text import MIMEText
import requests
from twilio.rest import Client

from pptx import Presentation
from pptx.util import Inches
from flask import send_file

from flask_cors import CORS
import schedule

from werkzeug.utils import secure_filename
import pandas as pd
from tempfile import TemporaryDirectory


UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'pdf', 'doc', 'docx'}

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.secret_key = os.environ.get('SECRET_KEY', '2309489358910358')
socketio = SocketIO(app)
CORS(app)
toastr = Toastr(app)


# Ensure the upload folder exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

DATA_DIR = 'data'
CHAT_LOG_FILE = 'chat_log.json'
config_path = 'c:\\temp\\config.json'
#config_path = '/tmp/config.json'

logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

def load_json(filename):
    logger.info('Dies ist ein Info-Log-Eintrag')
    with open(os.path.join(DATA_DIR, filename), 'r') as file:
        return json.load(file)
    

def save_json(filename, data):
    with open(os.path.join(DATA_DIR, filename), 'w') as file:
        json.dump(data, file, indent=4)

@app.route('/')
def dashboard():
    notifications = load_json('notifications.json')
    exercises = load_json('exercises.json')
    return render_template('dashboard.html', notifications=notifications, exercises=exercises)




@app.route('/upload_file', methods=['POST'])
def upload_file():
    if 'file' not in request.files or 'title' not in request.form or 'inject_id' not in request.form:
        return jsonify({'status': 'error', 'message': 'Missing file, title, or inject_id'})

    file = request.files['file']
    title = request.form['title']
    inject_id = int(request.form['inject_id'])

    if file.filename == '':
        return jsonify({'status': 'error', 'message': 'No selected file'})
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
        file_url = url_for('uploaded_file', filename=filename)

        injects = load_json('injects.json')
        for inject in injects:
            if inject['id'] == inject_id:
                if 'files' not in inject:
                    inject['files'] = []
                inject['files'].append({'title': title, 'url': file_url})
                save_json('injects.json', injects)
                break

        return jsonify({'status': 'success', 'file_url': file_url})
    else:
        return jsonify({'status': 'error', 'message': 'File type not allowed'})


@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)


###############################
### C O M U N I C A T I O N ###
###############################

@app.route('/communication')
def communication():
    return render_template('communication.html')


def save_message(message):
    try:
        with open(CHAT_LOG_FILE, 'r') as file:
            chat_log = json.load(file)
    except FileNotFoundError:
        chat_log = []

    message['timestamp'] = datetime.now().isoformat()
    chat_log.append(message)

    with open(CHAT_LOG_FILE, 'w') as file:
        json.dump(chat_log, file, indent=4)



@app.route('/get_chat_log')
def get_chat_log():
    try:
        with open(CHAT_LOG_FILE, 'r') as file:
            chat_log = json.load(file)
    except FileNotFoundError:
        chat_log = []

    return jsonify(chat_log)



class ChatNamespace(Namespace):
    def on_connect(self):
        logging.info("Client connected to chat namespace")

    def on_disconnect(self):
        logging.info("Client disconnected from chat namespace")

    def on_send_message(self, data):
        save_message(data)
        self.emit('message', data, namespace='/chat')

socketio.on_namespace(ChatNamespace('/chat'))



@app.route('/')
def index():
    return render_template('index.html')

@socketio.on('send_message')
def handle_send_message_event(data):
    save_message(data)
    socketio.emit('message', data, to=None)


@app.route('/reports')
def reports():
    reports = load_json('reports.json')
    return render_template('reports.html', reports=reports)

@app.route('/injects', endpoint='injects')
def injects_route():
    injects_data = load_json('injects.json')
    return render_template('injects.html', injects=injects_data)

def get_inject_by_id(inject_id):
    injects_data = load_json('injects.json')
    for inject in injects_data:
        if inject['id'] == inject_id:
            return inject
    return None

@app.template_filter('id_to_name')
def id_to_name(id):
    scenarios = load_json('scenarios.json')  # Laden Sie die Szenarien aus der JSON-Datei
    scenario = next((s for s in scenarios if s['id'] == str(id)), None)
    if scenario is not None:
        return scenario['name']
    else:
        return 'Unbekanntes Szenario'

###########################
#### S C E N A R I O S ####
###########################


@app.route('/scenarios')
def scenarios():
    try:
        logging.info('Loading scenarios from scenarios.json')
        scenarios = load_json('scenarios.json')
        logging.info('Loaded scenarios: %s', scenarios)
        return render_template('scenarios.html', scenarios=scenarios)
    except Exception as e:
        logging.error('Failed to load scenarios: %s', e)
        return "Failed to load scenarios", 500

@app.route('/new_scenario', methods=['GET', 'POST'])
def new_scenario():
    try:
        if request.method == 'POST':
            data = request.get_json()
            scenarios = load_json('scenarios.json')
            new_scenario = {
                'id': max(s['id'] for s in scenarios) + 1 if scenarios else 1,
                'name': data['name'],
                'difficulty': data['difficulty'],
                'description': data['description']
            }
            scenarios.append(new_scenario)
            save_json('scenarios.json', scenarios)
            return jsonify({'status': 'success'})
        return render_template('new_scenario.html')
    except Exception as e:
        logging.error(f"Failed to create scenario: {e}")
        return "Failed to create scenario", 500


@app.route('/edit_scenario/<int:id>', methods=['GET', 'POST'])
def edit_scenario(id):
    try:
        scenario = load_scenario(id)
        if request.method == 'POST':
            data = request.get_json()
            scenario.update({
                'name': data['name'],
                'difficulty': data['difficulty'],
                'description': data['description']
            })
            save_scenario(scenario)
            return jsonify({'status': 'success'})
        return render_template('edit_scenario.html', scenario=scenario)
    except Exception as e:
        logging.error(f"Failed to load scenario: {e}")
        return "Failed to load scenario", 500

def save_scenario(updated_scenario):
    scenarios = load_json('scenarios.json')
    for i, scenario in enumerate(scenarios):
        if scenario['id'] == updated_scenario['id']:
            scenarios[i] = updated_scenario
            break
    save_json('scenarios.json', scenarios)



def load_scenario(id):
    logging.debug('load_scenario called with id: %s', id)
    scenarios = load_json('scenarios.json')  # Load the exercises from the JSON file
    logging.debug('Loaded scenarios: %s', scenarios)
    for scenario in scenarios:
        if scenario['id'] == int(id):  # Convert the ID to an integer before comparing
            logging.debug('Found matching scenario: %s', scenario)
            return scenario
    logging.debug('No matching scenario found')
    return None


###########################
#### E X E R C I S E S ####
###########################


@app.route('/schedule_exercise', methods=['POST'])
def schedule_exercise():
    data = request.json
    logging.info("[SCHEX.01] JSON loaded from request body")
    if data is None:
        logging.info("Missing JSON in request body")
        return jsonify({"status": "error", "message": "Missing JSON in request body"}), 400

    exercises = load_json('exercises.json')
    logging.info("[SCHEX.02] Exercises loaded")

    # Generate a new unique ID for the exercise
    if exercises:
        new_id = max(exercise['id'] for exercise in exercises) + 1
    else:
        new_id = 1

    data['id'] = new_id
    logging.info("[SCHEX.03] Data updated with new ID")

    # Create a new ordered dictionary with 'id' first
    ordered_data = {'id': data['id']}
    for key, value in data.items():
        if key != 'id':
            ordered_data[key] = value

    exercises.append(data)
    save_json('exercises.json', exercises)
    logging.info(f"Exercise with id {new_id} has been scheduled")
    if data['type'] == 'timed':
        schedule_timed_exercise(data)
    return jsonify({"status": "scheduled", "id": new_id}), 200

def schedule_timed_exercise(exercise):
    injects = load_json('injects.json')
    for idx, inject_idx in enumerate(exercise['inject_order']):
        inject = injects[inject_idx]
        delay = idx * exercise['interval'] * 60  # convert minutes to seconds
        Timer(delay, execute_inject, [inject]).start()
        logging.info(f"Inject {inject_idx} has been scheduled with a delay of {delay} seconds")



@app.route('/update_exercise/<int:id>', methods=['POST'])
def update_exercise(id):
    logging.info('[UE.1] Updating exercise with ID: %s', id)

    exercises = load_json('exercises.json')
    logging.info('[UE.2] Loaded exercises from JSON file')

    for exercise in exercises:
        if exercise['id'] == int(id):
            logging.info('[UE.3] Found exercise with matching ID')

            exercise['name'] = request.form['name']
            exercise['description'] = request.form['description']
            exercise['target_team'] = request.form['target_team']
            exercise['planned'] = request.form['planned']
            exercise['execution_type'] = request.form['execution_type']

            inject_order = request.form['inject_order'].split(',')
            exercise['inject_order'] = [int(i) for i in inject_order if i]

            logging.info('[UE.4] Updated exercise properties')

            save_json('exercises.json', exercises)
            logging.info('[UE.5] Saved exercises to JSON file')

            flash('Exercise updated successfully')
            logging.info('[UE.6] Flashed message: Exercise updated successfully')

            return redirect(url_for('exercises'))

    logging.warning('[UE.7] No exercise found with ID: %s', id)
    flash('Exercise not found')
    return redirect(url_for('exercises'))


@app.route('/schedule_exercise_form')
def schedule_exercise_form():
    injects = load_json('injects.json')
    return render_template('exercise_planung.html', injects=injects)

@app.route('/exercises')
def exercises():
    exercises = load_json('exercises.json')  # Load the exercises from the JSON file
    #print(exercises)  # Print the exercises data
    return render_template('exercises.html', exercises=exercises)  # Pass the exercises to the template


def load_exercise(id):
    exercises = load_json('exercises.json')  # Load the exercises from the JSON file
    for exercise in exercises:
        if exercise['id'] == int(id):  # Convert the ID to an integer before comparing
            return exercise
    return None

@app.route('/edit_exercise/<id>')
def edit_exercise(id):
    exercise = load_exercise(id)
    injects = load_json('injects.json')  # Load injects from JSON file
    if exercise is not None:
        return render_template('edit_exercise.html', exercise=exercise, injects=injects)
    else:
        return "Exercise not found", 404



@socketio.on('start_exercise')
def perform_exercise(exercise):
    logging.info("[PE.01] perform_exercise called")
    socketio.emit('message', {'data': 'Execution started!'})

    if exercise is None:
        logging.error("[*] Error: exercise cannot be None")
        socketio.emit('message', {'data': 'Error: exercise not found'})
        return False

    logging.info("[PE.02] Loading exercises from exercises.json")
    exercises = load_json('exercises.json')

    logging.info("[PE.03] Loading injects from injects.json")
    injects = load_json('injects.json')

    if not 'inject_order' in exercise:
        logging.error("[*] No inject_order in exercise")
        socketio.emit('message', {'data': 'Error: No inject_order in exercise'})
        return False

    logging.info("[PE.04] Running through injects in the order specified in the exercise")
    for inject_id in exercise['inject_order']:
        logging.info(f"[PE.05] Looking for inject with ID {inject_id}")
        inject = next((inject for inject in injects if inject['id'] == inject_id), None)

        if inject is None:
            logging.error(f"[*] No inject found with ID {inject_id}")
            continue

        title = inject.get('title', 'No title')
        duration = inject.get('duration', 10)  # Default to 10 seconds if duration is not specified
        # Get the communication_type
        communication_type = inject.get('communication_type')
        # Get the assigned scenarios
        # Pause for 5 seconds
        
        for i in range(5, 0, -1):
            socketio.emit('message', {'data': f'Countdown: {i}'})
            time.sleep(1)
        
        logging.info(f"[PE.05] Executing inject {inject_id}: {title} for {duration} seconds")
        socketio.emit('message', {'data': f'Inject {inject_id} wird ausgeführt: {title} für {duration} Sekunden per {communication_type}'})

        # Check if communication_type is not None
        if communication_type:
            # Call the appropriate function based on the communication_type
            if communication_type == 'text':
                #textnote(f"[{exercise.get('name', 'No title')}] {inject.get('title', 'No title')}", inject.get('nachrichtentext', 'No description'))
                #whatsapp_notification()
                message = "[" + exercise.get('name', 'No title') + "] " + inject.get('title', 'No title') + ", " + inject.get('nachrichtentext', 'No description')
                textnote_internal("forensician", message)


            elif communication_type == 'email':
                email(f"[{exercise.get('name', 'No title')}] {inject.get('title', 'No title')}", inject.get('nachrichtentext', 'No description'))
                #make_call_route()

        logging.info(f"[PE.06] Waiting for {duration} seconds")
        time.sleep(duration)

    logging.info("[*] Finished executing injects")
    socketio.emit('message', {'data': 'Exercise finished'})

    # Update the last_performed field with the current date and time
    exercise['last_performed'] = datetime.now().isoformat()

    # Find the exercise in the list of exercises and update it
    for i, ex in enumerate(exercises):
        if ex['id'] == exercise['id']:
            exercises[i] = exercise
            break

    # Save the updated exercises back to the JSON file
    save_json('exercises.json', exercises)

    socketio.emit('message', {'data': 'Scenario updated successfully'})

    return True



@app.route('/execute_exercise/<int:exercise_id>', methods=['POST'])
def execute_exercise(exercise_id):
    logging.info(f"[EE.01] execute_exercise called with exercise_id: {exercise_id}")

    logging.info("[EE.02] Loading exercise")
    exercise = load_exercise(exercise_id)

    # Check if the exercise exists
    if exercise is None:
        logging.error("[*] Error: Exercise not found")
        return render_template('error.html', message="Exercise not found"), 404


    logging.info("[EE.03] Starting a new thread to perform the exercise")
    # Perform the exercise in a separate thread
    thread = Thread(target=perform_exercise, args=(exercise,))
    thread.start()

    logging.info("[EE.04] Rendering the executed_exercise.html page")
    # Render the executed_exercise.html page immediately
    return render_template('executed_exercise.html', exercise=exercise, status="Exercise is being executed")


@app.route('/create_and_download_presentation_for_exercise/<int:exercise_id>', methods=['POST'])
def create_and_download_presentation_for_exercise(exercise_id):
    exercises = load_json('exercises.json')
    injects = load_json('injects.json')
    exercise = next((ex for ex in exercises if ex['id'] == exercise_id), None)
    
    if exercise:
        file_path = create_powerpoint_for_exercise(exercise, injects)
        return send_file(file_path, as_attachment=True)
    else:
        flash('Exercise not found')
        return redirect(url_for('exercises'))

def create_powerpoint_for_exercise(exercise, injects):
    inject_map = {inject['id']: inject for inject in injects}

    # Load the PowerPoint template
    template_path = 'static/template.pptx'  # Path to your PowerPoint template
    prs = Presentation(template_path)

    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = exercise.get('name', 'No title')
    subtitle.text = exercise.get('description', 'No description')

    # Create slides for each inject
    inject_order = exercise.get('inject_order', [])
    comments = exercise.get('inject_comment', {})

    for inject_id in inject_order:
        inject = inject_map.get(inject_id, {})
        inject_title = inject.get('title', 'No title')
        inject_duration = inject.get('duration', 'No duration')
        inject_comment = comments.get(str(inject_id), 'No comment')
        inject_nachrichtentextPlain = inject.get('nachrichtentextPlain', 'No description')

        slide_layout = prs.slide_layouts[1]  # You can use a different layout if needed
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = f"Inject {inject_id}: {inject_title}"
        #content.text = f"Duration: {inject_duration} seconds\n\nMessage: {inject_nachrichtentextPlain}\n\nComment: {inject_comment}"
        content.text = f"{inject_nachrichtentextPlain}"

    file_path = f'static/{exercise["name"].replace(" ", "_")}_presentation.pptx'
    prs.save(file_path)
    return file_path



# Fügen Sie diese Zeile am Anfang Ihrer Datei hinzu
current_inject_indices = {}


@app.route('/next_inject/<int:exercise_id>')
def next_inject(exercise_id):
    logging.info('[NextI] Fetching next inject for exercise_id: %s', exercise_id)
    exercise = load_exercise(exercise_id)
    injects = load_json('injects.json')
    if not exercise or not exercise.get('inject_order'):
        logging.info('[NextI] No injects found for exercise_id: %s', exercise_id)
        return jsonify({"status": "error", "message": "No injects found"}), 404

    # Holen Sie den aktuellen Index für diese exerciseId, oder setzen Sie ihn auf 0, wenn er noch nicht existiert
    current_index = current_inject_indices.get(exercise_id, 0)
    logging.info('[NextI] Current index for exercise_id %s: %s', exercise_id, current_index)

    # Überprüfen Sie, ob es noch weitere Injects gibt
    if current_index >= len(exercise['inject_order']):
        logging.info('[NextI] No more injects for exercise_id: %s', exercise_id)
        return jsonify({"status": "error", "message": "No more injects"}), 404

    next_inject_id = exercise['inject_order'][current_index]
    next_inject = next((inject for inject in injects if inject['id'] == next_inject_id), None)
    if not next_inject:
        logging.info('[NextI] Next inject not found for exercise_id: %s', exercise_id)
        return jsonify({"status": "error", "message": "Next inject not found"}), 404

    # Aktualisieren Sie den Index für die nächste Anforderung
    current_inject_indices[exercise_id] = current_index + 1
    logging.info('[NextI] Updated index for exercise_id %s: %s', exercise_id, current_index + 1)

    return jsonify({"status": "success", "inject": next_inject}), 200



@app.route('/update_comment', methods=['POST'])
def update_comment():
    logging.info('[UC.1] update_comment function called')

    try:
        data = request.get_json()
        logging.info(f'[UC.2] Received data: {data}')

        if 'id' in data and 'index' in data and 'comment' in data:
            try:
                exercise_id = int(data['id'])
                inject_index = int(data['index'])
                logging.info(f'[UC.3] Converting id and index to integer, exercise_id: {exercise_id}, inject_index: {inject_index}')
            except ValueError:
                logging.error('[*] Error: Invalid type for id or index')
                return jsonify({'status': 'error', 'message': 'Invalid type for id or index'})

            comment = data['comment']
            logging.info(f'[UC.4] Exercise ID: {exercise_id}, Inject Index: {inject_index}, Comment: {comment}')

            try:
                exercises = load_json('exercises.json')
                logging.info(f'[UC.5] Loaded exercises: {exercises}')
            except FileNotFoundError:
                logging.error('[*] Error: exercises.json file not found')
                return jsonify({'status': 'error', 'message': 'File not found'})
            except json.JSONDecodeError:
                logging.error('[*] Error: JSON decode error')
                return jsonify({'status': 'error', 'message': 'Error reading JSON file'})

            # Suche nach der Übung mit der entsprechenden ID
            exercise = None
            for ex in exercises:
                if ex['id'] == exercise_id:
                    exercise = ex
                    break

            if exercise:
                if 'inject_comment' not in exercise or not isinstance(exercise['inject_comment'], dict):
                    exercise['inject_comment'] = {}
                elif isinstance(exercise['inject_comment'], list):
                    # Konvertiere die Liste zu einem Wörterbuch
                    inject_comment_dict = {}
                    for i, val in enumerate(exercise['inject_comment']):
                        inject_comment_dict[str(i)] = val
                    exercise['inject_comment'] = inject_comment_dict
                
                exercise['inject_comment'][str(inject_index)] = comment
                logging.info(f'[UC.6] Updated inject_comment: {exercise["inject_comment"]}')
                logging.info(f'[UC.7] Updated exercise object: {exercise}')

                # Speichern der gesamten Liste zurück in die JSON-Datei
                try:
                    save_json('exercises.json', exercises)
                except IOError:
                    logging.error('[*] Error: Unable to write to exercises.json file')
                    return jsonify({'status': 'error', 'message': 'Unable to write to file'})

                logging.info('[UC.8] Saved exercises back to JSON file')

                return jsonify({'status': 'success'})
            else:
                logging.error(f'[*] Error: Exercise not found. Received id: {exercise_id}')
                return jsonify({'status': 'error', 'message': 'Exercise not found'})
        else:
            logging.error(f'[*] Error: Missing key in data. Received data: {data}')
            return jsonify({'status': 'error', 'message': 'Missing key in data'})

    except Exception as e:
        logging.error(f'[*] Error: {str(e)}')
        return jsonify({'status': 'error', 'message': 'Internal server error'})
    



@app.route('/view_inject_comments/<int:exercise_id>', methods=['GET'])
def view_inject_comments(exercise_id):
    logging.info(f"[VIC.01] view_inject_comments called with exercise_id: {exercise_id}")

    logging.info("[VIC.02] Loading exercise")
    exercise = load_exercise(exercise_id)

    # Check if the exercise exists
    if exercise is None:
        logging.error("[*] Error: Exercise not found")
        return render_template('error.html', message="Exercise not found"), 404

    # Check if 'inject_comments' exists in exercise
    if 'inject_comments' in exercise:
        print(exercise['inject_comment'])  # Debug output
    else:
        logging.error("[*] Error: 'inject_comments' not found in exercise")

    logging.info("[VIC.03] Rendering the inject_comments.html page")
    # Render the inject_comments.html page
    return render_template('inject_comments.html', exercise=exercise)



#### I N J E C T S ####

@app.route('/edit_inject_form/<int:inject_id>')
def edit_inject_form(inject_id):
    inject = get_inject_by_id(inject_id)
    if inject is None:
        return "Inject not found", 404
    scenarios = load_json('scenarios.json')
    return render_template('edit_inject.html', inject=inject, scenarios=scenarios)

@app.route('/execute_inject_form')
def execute_inject_form():
    return render_template('execute_exercise.html')


@app.route('/execute_inject/<exercise_name>/<inject_idx>')
def execute_inject_route(exercise_name, inject_idx):
    exercises = load_json('exercises.json')
    injects = load_json('injects.json')
    exercise = next((ex for ex in exercises if ex['name'] == exercise_name), None)
    if exercise:
        inject = injects[int(inject_idx)]
        execute_inject(inject)
    return jsonify({"status": "executed"}), 200

@app.route('/execute_inject/<int:inject_id>', methods=['POST'])
def execute_single_inject(inject_id):
    logging.info(f"[EI.01] execute_single_inject called with inject_id: {inject_id}")

    injects = load_json('injects.json')
    inject = next((inject for inject in injects if inject['id'] == inject_id), None)

    if inject is None:
        logging.error("[*] Error: Inject not found")
        return jsonify({"status": "error", "message": "Inject not found"}), 404

    # Execute the inject (placeholder for the actual logic)
    execute_inject(inject)

    logging.info(f"[EI.02] Inject {inject_id} executed")
    return jsonify({"status": "success", "message": f"Inject {inject_id} executed"}), 200

def execute_inject(inject):
    logging.info(f"Executing inject: {inject['title']}")
    # Simulate the inject execution logic
    communication_type = inject.get('communication_type')
    if communication_type == 'text':
        message = inject.get('nachrichtentext', 'No description')
        textnote_internal("forensician", message)
    elif communication_type == 'email':
        email(inject.get('title', 'No title'), inject.get('nachrichtentext', 'No description'))
    # Add other communication types as needed
    logging.info(f"Inject executed: {inject['title']} via {communication_type}")


@app.route('/new_inject', methods=['GET', 'POST'])
def new_inject():
    try:
        scenarios = load_json('scenarios.json')
        injects = load_json('injects.json')

        if request.method == 'POST':
            data = request.get_json()

            print("Received data:", data)
            logging.info(f"Nachrichtentext as received from client: {data.get('nachrichtentext')}")

            # Set default values for missing fields
            new_inject = {
                'id': len(injects) + 1,  # Assuming id is a sequential integer
                'title': data.get('title', ""),
                'description': data.get('description', ""),
                'exercise_benefit': data.get('exercise_benefit', ""),
                'expected_response': data.get('expected_response', ""),
                'communication_type': data.get('communication_type', ""),
                'duration': data.get('duration', ""),
                'nachrichtentext': data.get('nachrichtentext', ""),
                'nachrichtentextPlain': data.get('nachrichtentextPlain', ""),
                'files': data.get('files', [])
            }

            injects.append(new_inject)
            save_json('injects.json', injects)

            logging.debug(f"Nachrichtentext as packed into dictionary for JSON: {new_inject['nachrichtentext']}")

        return render_template('new_inject.html', scenarios=scenarios)

    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return "An error occurred while processing your request", 500



@app.route('/edit_inject/<int:inject_id>', methods=['GET', 'POST'])
def edit_inject(inject_id):
    try:
        inject = get_inject_by_id(inject_id)
        scenarios = load_json('scenarios.json')
        injects = load_json('injects.json')

        if inject is None:
            return "Inject not found", 404

        if request.method == 'POST':
            data = request.get_json()

            print("Received data:", data)
            logging.info(f"Nachrichtentext as received from client: {data.get('nachrichtentext')}")

            # Set default values for missing fields
            for field in ['title', 'description', 'exercise_benefit', 'expected_response', 'communication_type', 'duration', 'nachrichtentext', 'nachrichtentextPlain']:
                inject.setdefault(field, "")

            updated_inject = {
                'id': inject_id,
                'title': data.get('title', inject['title']),
                'description': data.get('description', inject['description']),
                'exercise_benefit': data.get('exercise_benefit', inject['exercise_benefit']),
                'expected_response': data.get('expected_response', inject['expected_response']),
                'communication_type': data.get('communication_type', inject['communication_type']),
                'duration': data.get('duration', inject['duration']),
                'nachrichtentext': data.get('nachrichtentext', inject['nachrichtentext']),
                'nachrichtentextPlain': data.get('nachrichtentextPlain', inject['nachrichtentextPlain']),
                'files': inject.get('files', [])
            }

            inject_index = next((index for (index, d) in enumerate(injects) if d['id'] == inject_id), None)

            if inject_index is not None:
                injects[inject_index] = updated_inject
            else:
                return "Inject not found in database", 404

            save_json('injects.json', injects)

            logging.debug(f"Nachrichtentext as packed into dictionary for JSON: {updated_inject['nachrichtentext']}")

        return render_template('edit_inject.html', inject=inject, scenarios=scenarios)

    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return "An error occurred while processing your request", 500


@app.route('/export_injects', methods=['GET'])
def export_injects():
    try:
        injects = load_json('injects.json')
        
        # Data preparation for Excel
        data = []
        for inject in injects:
            assigned_scenarios = inject.get('assigned_scenarios', [])
            data.append({
                'Id': inject['id'],
                'Title': inject['title'],
                'Type of communication': inject['communication_type'],
                'Description': inject['description'],
                'Exercise benefit': inject['exercise_benefit'],
                'Expected handling': inject['expected_response'],
                'Assigned scenarios': ', '.join([id_to_name(scenario_id) for scenario_id in assigned_scenarios])
            })
        
        # Convert to DataFrame
        df = pd.DataFrame(data)
        
        with TemporaryDirectory() as tempdir:
            excel_path = os.path.join(tempdir, 'injects_export.xlsx')
            df.to_excel(excel_path, index=False)
            
            # Read the file content and send it as a response
            with open(excel_path, 'rb') as f:
                file_content = f.read()

        return send_file(
            io.BytesIO(file_content),
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name='injects_export.xlsx'  # Use download_name instead of attachment_filename
        )
    except Exception as e:
        logging.error(f"An error occurred while exporting injects: {e}")
        return jsonify({'message': 'An error occurred while exporting injects', 'error': str(e)}), 500




##### N O T E S #####


def textnote_internal(user, message):
    data = {'user': user, 'message': message}
    save_message(data)
    socketio.emit('message', data, namespace='/chat')

def textnote(title, description):
    # Load the config file
    #config_path = 'c:\\temp\\config.json'
    #config_path = '\\tmp\\config.json'
    try:
        with open(config_path, 'r') as f:
            config = json.load(f)
    except FileNotFoundError:
        logging.error(f"[textnote] Config file not found at {config_path}")
        return
    except json.JSONDecodeError:
        logging.error(f"[textnote] Failed to parse JSON from config file at {config_path}")
        return

    # Define the Slack API URL and headers
    url = "https://slack.com/api/chat.postMessage"
    headers = {
        "Authorization": f"Bearer {config.get('text_token_keys')}",
        "Content-type": "application/json"
    }
    print(f"[**********] Bearer {config.get('text_token_keys')}")
    # Define the data to send
    data = {
        "channel": "alerts",
        "text": f"{title}: {description}"
    }

def threema_message(title, description):
    # Load the config file
    #config_path = 'c:\\temp\\config.json'
    #config_path = '\\tmp\\config.json'
    try:
        with open(config_path, 'r') as f:
            config = json.load(f)
    except FileNotFoundError:
        logging.error(f"[threema_message] Config file not found at {config_path}")
        return
    except json.JSONDecodeError:
        logging.error(f"[threema_message] Failed to parse JSON from config file at {config_path}")
        return

    # Define the Threema Gateway API URL and headers
    url = "https://msgapi.threema.ch/send_simple"
    headers = {
        "Content-type": "application/json"
    }

    # Define the data to send
    data = {
        "from": config.get('threema_from'),
        "to": config.get('threema_to'),
        "text": f"{title}: {description}"
    }

    # Send the request
    response = requests.post(url, headers=headers, data=json.dumps(data), auth=(config.get('threema_api_identity'), config.get('threema_api_secret')))

    # Check the response
    if response.status_code != 200:
        logging.error(f"[threema_message] Failed to send message: {response.content}") 

    # Send the request
    try:
        response = requests.post(url, headers=headers, data=json.dumps(data))
        response.raise_for_status()
        logging.info(f"[textnote] Sent Slack notification with title: {title} and description: {description}")
    except requests.exceptions.RequestException as e:
        logging.error(f"[textnote] Failed to send Slack notification. RequestException: {e}")
    except Exception as e:
        logging.error(f"[textnote] Failed to send Slack notification with title: {title} and description: {description}. Error: {e}")


def email(title, description):
    # Load the config file
    #config_path = 'c:\\temp\\config.json'
    #config_path = '\\tmp\\config.json'
    with open(config_path, 'r') as f:
        config = json.load(f)

    # Define the SMTP server parameters
    server = config.get('email_server')
    port = 25  # Change this to 25
    username = config.get('email_user')
    password = config.get('email_password')

    # Define the email parameters
    sender = config.get('email_from')
    receiver = config.get('email_to')
    subject = title
    body = description

    # Create the email
    msg = MIMEText(body, 'html')
    msg['Subject'] = subject
    msg['From'] = sender
    msg['To'] = receiver

    # Send the email
    try:
        smtp_server = smtplib.SMTP(server, port)
        smtp_server.ehlo()
        smtp_server.starttls()  # Call this before login
        smtp_server.login(username, password)
        smtp_server.sendmail(sender, receiver, msg.as_string())
        smtp_server.close()
        #logging.debug(f"[emailnote] Sent email notification with title: {title} and description: {description}")
    except Exception as e:
        logging.error(f"[emailnote] Failed to send email notification with title: {title} and description: {description}. Error: {e}")



def gmail(title, description):

    with open(config_path, 'r') as f:
        config = json.load(f)

    # Define the SMTP server parameters
    server = 'smtp.gmail.com'  # Google Mail SMTP server address
    port = 587  # Google Mail SMTP server port
    username = config.get('email_user')
    password = config.get('email_password')

    # Define the email parameters
    sender = config.get('email_from')
    receiver = config.get('email_to')
    subject = title
    body = description

    # Create the email
    msg = MIMEText(body)
    msg['Subject'] = subject
    msg['From'] = sender
    msg['To'] = receiver

    # Send the email
    try:
        smtp_server = smtplib.SMTP(server, port)
        smtp_server.ehlo()
        smtp_server.starttls()  # Call this before login
        smtp_server.login(username, password)
        smtp_server.sendmail(sender, receiver, msg.as_string())
        smtp_server.close()
        logging.info(f"[emailnote] Sent email notification with title: {title} and description: {description}")
    except Exception as e:
        logging.error(f"[emailnote] Failed to send email notification with title: {title} and description: {description}. Error: {e}")


def make_call_route():
    # Load the config file
    #config_path = 'c:\\temp\\config.json'
    with open(config_path, 'r') as f:
        config = json.load(f)

    account_sid = config.get('twilio_account_sid')
    auth_token = config.get('twilio_auth_token')
    client = Client(account_sid, auth_token)

    call = client.calls.create(
        twiml='<Response><Say>Ahoy, World!</Say></Response>',
        to=config.get('twilio_to'),
        from_=config.get('twilio_from')
                    )

    print(call.sid)


def whatsapp_notification(title, description):
        # Load the config file
    #config_path = 'c:\\temp\\config.json'
    #config_path = '\\tmp\\config.json'
    with open(config_path, 'r') as f:
        config = json.load(f)

    account_sid = config.get('twilio_account_sid')
    auth_token = config.get('twilio_auth_token')
    client = Client(account_sid, auth_token)

    message = client.messages.create(
    from_='whatsapp:+14155238886',
    body='{title}: {description}',
    #to='whatsapp:+491726719759'
    to='whatsapp:+491743039038'
)



@app.route('/save_exercise_order', methods=['POST'])
def save_exercise_order():
    data = request.json
    exercises = load_json('exercises.json')
    exercise = next((ex for ex in exercises if ex['name'] == data['exercise_name']), None)
    if exercise:
        exercise['inject_order'] = data['order']
        save_json('exercises.json', exercises)
    return jsonify({"status": "success"}), 200

#### A P I ####

@app.route('/api/notifications', methods=['GET', 'POST'])
def api_notifications():
    if request.method == 'GET':
        return jsonify(load_json('notifications.json'))
    elif request.method == 'POST':
        data = request.json
        save_json('notifications.json', data)
        return jsonify({"status": "success"}), 200

@app.route('/api/exercises', methods=['GET', 'POST'])
def api_exercises():
    if request.method == 'GET':
        return jsonify(load_json('exercises.json'))
    elif request.method == 'POST':
        data = request.json
        save_json('exercises.json', data)
        return jsonify({"status": "success"}), 200

@app.route('/api/scenarios', methods=['GET', 'POST'])
def api_scenarios():
    if request.method == 'GET':
        return jsonify(load_json('scenarios.json'))
    elif request.method == 'POST':
        data = request.json
        save_json('scenarios.json', data)
        return jsonify({"status": "success"}), 200

@app.route('/api/reports', methods=['GET', 'POST'])
def api_reports():
    if request.method == 'GET':
        return jsonify(load_json('reports.json'))
    elif request.method == 'POST':
        data = request.json
        save_json('reports.json', data)
        return jsonify({"status": "success"}), 200

@app.route('/api/injects', methods=['GET', 'POST'])
def api_injects():
    if request.method == 'GET':
        return jsonify(load_json('injects.json'))
    elif request.method == 'POST':
        data = request.json
        save_json('injects.json', data)
        return jsonify({"status": "success"}), 200




@app.route('/edit_config', methods=['GET', 'POST'])
def edit_config():
    #config_path = 'c:\\temp\\config.json'
    #config_path = '\\tmp\\config.json'
    if request.method == 'POST':
        # Load the current config data
        with open(config_path, 'r') as f:
            config = json.load(f)
        # Update the config data with the form data
        form_data = request.form.to_dict()
        config.update(form_data)
        # Write the updated config data back to the file
        with open(config_path, 'w') as f:
            json.dump(config, f)
        # Redirect back to the config page
        return redirect(url_for('edit_config'))
    else:
        # Render the form with the current config data
        with open(config_path, 'r') as f:
            config = json.load(f)
        return render_template('config_form.html', config=config)

if __name__ == '__main__':
    #app.run(debug=True, port=5010)
    socketio.run(app, debug=True, host="0.0.0.0", port=5010)

